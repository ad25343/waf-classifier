"""
Generate WAF Classifier capabilities deck — light theme.
Run: python3 make_deck.py
Output: static/WAF_Classifier_Capabilities.pptx

Structure (13 slides):
  MAIN (5):  Title · Problem · 8 Modules · What Leadership Gets · Where AI is Used
  APPENDIX (8): Divider · Classification · Analytics · Merge · Teams · Quality · Lineage+Disputes · Settings
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Light Palette ─────────────────────────────────────────
BG          = RGBColor(0xF7, 0xF8, 0xFC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CARD        = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE     = RGBColor(0xF0, 0xF2, 0xF9)
BORDER      = RGBColor(0xD8, 0xDC, 0xEE)
ACCENT      = RGBColor(0x5B, 0x52, 0xE8)   # indigo
ACCENT2     = RGBColor(0x00, 0xA8, 0x90)   # teal
ACCENT_PALE = RGBColor(0xEE, 0xED, 0xFF)
TEAL_PALE   = RGBColor(0xDF, 0xF6, 0xF3)
HEAD        = RGBColor(0x0F, 0x11, 0x1A)   # near-black heading
TEXT        = RGBColor(0x1E, 0x22, 0x3A)
MUTED       = RGBColor(0x6B, 0x73, 0x96)
LLM_BG      = RGBColor(0xED, 0xEB, 0xFF)
LLM_FG      = RGBColor(0x4A, 0x40, 0xCC)
NOLLM_BG    = RGBColor(0xD9, 0xF5, 0xF0)
NOLLM_FG    = RGBColor(0x00, 0x87, 0x72)
WARN        = RGBColor(0xD9, 0x7A, 0x06)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
RED         = RGBColor(0xDC, 0x26, 0x26)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

TOTAL = 13


# ── Helpers ───────────────────────────────────────────────

def new_slide():
    sl = prs.slides.add_slide(blank_layout)
    bg = sl.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return sl


def rect(sl, x, y, w, h, fill=CARD, lc=None, lw=Pt(0.75)):
    shp = sl.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if lc:
        shp.line.color.rgb = lc; shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp


def rrect(sl, x, y, w, h, fill=CARD, lc=BORDER, lw=Pt(0.75)):
    shp = sl.shapes.add_shape(5, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = lc; shp.line.width = lw
    shp.adjustments[0] = 0.04
    return shp


def txt(sl, text, x, y, w, h, size=12, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False):
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    r.font.italic = italic
    return txb


def footer(sl, cur):
    rect(sl, 0, SLIDE_H - Inches(0.26), SLIDE_W, Inches(0.26), fill=WHITE, lc=BORDER, lw=Pt(0.4))
    pw = int(SLIDE_W * cur / TOTAL)
    if pw > 0:
        r = rect(sl, 0, SLIDE_H - Pt(3), pw, Pt(3), fill=ACCENT)
        r.line.fill.background()
    txt(sl, "WAF Classifier  ·  Capabilities Overview",
        Inches(0.3), SLIDE_H - Inches(0.25), Inches(5), Inches(0.24),
        size=8, color=MUTED)
    txt(sl, f"{cur} / {TOTAL}",
        SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.25), Inches(0.7), Inches(0.24),
        size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def slide_eyebrow(sl, text, x=Inches(0.55), y=Inches(0.22)):
    txt(sl, text.upper(), x, y, Inches(5), Pt(14), size=8.5, bold=True, color=ACCENT)


def slide_title(sl, text, x=Inches(0.55), y=Inches(0.42), w=Inches(12)):
    txt(sl, text, x, y, w, Inches(0.65), size=28, bold=True, color=HEAD)


def slide_sub(sl, text, x=Inches(0.55), y=Inches(1.02), w=Inches(10.5)):
    txt(sl, text, x, y, w, Inches(0.44), size=12.5, color=MUTED)


def divider(sl, y=Inches(1.0)):
    rect(sl, Inches(0.55), y, SLIDE_W - Inches(1.1), Pt(1), fill=BORDER)


def badge(sl, x, y, llm=False):
    bw = Inches(0.78) if llm else Inches(0.88)
    bg_c = LLM_BG if llm else NOLLM_BG
    fg_c = LLM_FG if llm else NOLLM_FG
    btxt = "🤖  LLM" if llm else "✓  No LLM"
    r = rrect(sl, x, y, bw, Pt(15), fill=bg_c, lc=bg_c)
    txt(sl, btxt, x, y - Pt(1), bw, Pt(17),
        size=7.5, bold=True, color=fg_c, align=PP_ALIGN.CENTER)


def cap_card(sl, x, y, w, h, icon, title, desc, llm=False):
    rrect(sl, x, y, w, h, fill=CARD, lc=BORDER)
    txt(sl, icon, x + Inches(0.16), y + Inches(0.1), Inches(0.38), Inches(0.38), size=15)
    badge(sl, x + w - (Inches(0.78) if llm else Inches(0.88)) - Inches(0.12),
          y + Inches(0.13), llm)
    txt(sl, title, x + Inches(0.16), y + Inches(0.42),
        w - Inches(0.32), Inches(0.28), size=11, bold=True, color=HEAD)
    txt(sl, desc, x + Inches(0.16), y + Inches(0.68),
        w - Inches(0.32), h - Inches(0.76), size=9.5, color=MUTED)


# ════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE  (MAIN)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()

txt(sl, "WAF CLASSIFIER", Inches(0.55), Inches(0.32), Inches(8), Inches(0.3),
    size=10, bold=True, color=ACCENT)

txt(sl, "Are we working on the right things?",
    Inches(0.55), Inches(0.66), Inches(12.5), Inches(1.0),
    size=40, bold=True, color=HEAD)

txt(sl, "We deliver thousands of stories every quarter. Leadership needs to know — with confidence — "
        "where engineering spend is going, whether it's aligned to what matters, and whether we're "
        "shrinking the maintenance burden or growing it.",
    Inches(0.55), Inches(1.78), Inches(12.0), Inches(0.96),
    size=14, color=MUTED)

divider(sl, Inches(2.95))

txt(sl, "THE QUESTIONS LEADERSHIP NEEDS ANSWERED",
    Inches(0.55), Inches(3.12), Inches(8), Inches(0.3),
    size=9, bold=True, color=ACCENT)

questions = [
    ("🎯", "Are we aligned?",
     "Is every story tied to a strategic priority — or are we drifting?"),
    ("⚖️", "What's our Run vs. Change split?",
     "How much capacity is going to keep-the-lights-on vs. value-generating work?"),
    ("📊", "Can we trust the numbers?",
     "Is the WAF data accurate, auditable, and consistent across every team?"),
]
qw = Inches(3.9)
for i, (ico, q, sub) in enumerate(questions):
    qx = Inches(0.55) + i * (qw + Inches(0.21))
    qy = Inches(3.46)
    rrect(sl, qx, qy, qw, Inches(1.5), fill=WHITE, lc=BORDER)
    rect(sl, qx, qy, qw, Pt(3), fill=ACCENT).line.fill.background()
    txt(sl, ico, qx + Inches(0.18), qy + Inches(0.14), Inches(0.4), Inches(0.4), size=16)
    txt(sl, q, qx + Inches(0.62), qy + Inches(0.18),
        qw - Inches(0.74), Inches(0.34), size=13, bold=True, color=HEAD)
    txt(sl, sub, qx + Inches(0.18), qy + Inches(0.66),
        qw - Inches(0.36), Inches(0.78), size=10.5, color=MUTED)

divider(sl, Inches(5.18))

txt(sl, "WAF CLASSIFIER  —  AN AI-POWERED PLATFORM",
    Inches(0.55), Inches(5.32), Inches(8), Inches(0.3),
    size=9, bold=True, color=ACCENT)

txt(sl, "Classifies, verifies, and analyses every story against the Work Allocation Framework",
    Inches(0.55), Inches(5.62), Inches(12), Inches(0.4),
    size=16, bold=True, color=HEAD)

outcomes = [
    ("Consistent classification at scale",
     "AI-assisted tagging eliminates manual variance. Every story is classified the same way, every time."),
    ("Run vs. Change visibility",
     "Track maintenance burden vs. value-generating work across every team, epic, and Team of Teams."),
    ("Auditable, trusted data",
     "Every classification is logged with reasoning, confidence, and the WAF version that produced it."),
]
ow = Inches(3.9)
for i, (ttl, desc) in enumerate(outcomes):
    ox = Inches(0.55) + i * (ow + Inches(0.21))
    oy = Inches(6.16)
    rrect(sl, ox, oy, ow, Inches(0.94), fill=ACCENT_PALE, lc=ACCENT_PALE)
    txt(sl, ttl, ox + Inches(0.16), oy + Inches(0.08),
        ow - Inches(0.32), Inches(0.3), size=10.5, bold=True, color=ACCENT)
    txt(sl, desc, ox + Inches(0.16), oy + Inches(0.38),
        ow - Inches(0.32), Inches(0.54), size=9.5, color=TEXT)

footer(sl, 1)


# ════════════════════════════════════════════════════════════
#  SLIDE 2 — THE PROBLEM  (MAIN)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=RED).line.fill.background()
slide_eyebrow(sl, "The Problem")
slide_title(sl, "We deliver thousands of stories every quarter.")
txt(sl, "But can we honestly say they're all aligned, accurately tracked, and tied to what matters most for the organisation?",
    Inches(0.55), Inches(1.06), Inches(11.5), Inches(0.5),
    size=13, color=MUTED)
divider(sl, Inches(1.52))

rrect(sl, Inches(0.55), Inches(1.64), Inches(12.2), Inches(0.76),
      fill=RGBColor(0xFF, 0xF4, 0xF4), lc=RGBColor(0xF8, 0xC4, 0xC4))
txt(sl, "❝  Are we working on the right things?  Are we aligned?  Are we focused on initiatives that drive real value?  ❞",
    Inches(0.75), Inches(1.74), Inches(11.8), Inches(0.52),
    size=12.5, italic=True, color=RED, align=PP_ALIGN.CENTER)

pain_points = [
    ("📉", "Inconsistent classification",
     "Hundreds of teams tag stories against the WAF independently. Without a shared baseline, the same story gets classified differently depending on who's doing the tagging."),
    ("⏱", "Doesn't scale manually",
     "Reviewing and tagging thousands of stories per PI is slow, tedious, and error-prone. Central teams become a bottleneck. Deadlines get missed."),
    ("📋", "Poor story quality",
     "Stories missing acceptance criteria, source data, or clear outcomes can't be reliably classified — or delivered. Quality problems cascade."),
    ("🔍", "No audit trail",
     "When a classification is challenged — 'why is this tagged ORANGE?' — there's no record of the reasoning, who tagged it, or what version of the WAF was active."),
    ("📊", "Leadership can't trust the data",
     "If the Run vs. Change split or strategic category breakdown is based on inconsistent manual tagging, the numbers reported to senior leadership are unreliable."),
    ("🎯", "Misalignment goes undetected",
     "Without visibility across teams and epics, work that drifts away from strategic priorities — SPEED, tech modernisation, customer value — goes unnoticed until it's too late."),
]

cw, ch = Inches(3.9), Inches(1.56)
for i, (ico, ttl, dsc) in enumerate(pain_points):
    cx = Inches(0.55) + (i % 3) * (cw + Inches(0.18))
    cy = Inches(2.56) + (i // 3) * (ch + Inches(0.14))
    rrect(sl, cx, cy, cw, ch, fill=WHITE, lc=BORDER)
    rect(sl, cx, cy, cw, Pt(3), fill=RED).line.fill.background()
    txt(sl, ico, cx + Inches(0.14), cy + Inches(0.1), Inches(0.36), Inches(0.36), size=14)
    txt(sl, ttl, cx + Inches(0.55), cy + Inches(0.12), cw - Inches(0.65), Inches(0.3),
        size=11, bold=True, color=HEAD)
    txt(sl, dsc, cx + Inches(0.14), cy + Inches(0.5), cw - Inches(0.26), Inches(0.92),
        size=9.5, color=MUTED)

footer(sl, 2)


# ════════════════════════════════════════════════════════════
#  SLIDE 3 — 8 MODULES OVERVIEW  (MAIN)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Platform")
slide_title(sl, "8 Modules — One Platform")
slide_sub(sl, "Surface area at a glance. Each module is detailed in the appendix.")
divider(sl, Inches(1.42))

modules_detail = [
    ("🤖", "Classification",    "Single-story chat classification and async bulk file verification using AI",     "7 capabilities"),
    ("📊", "Analytics & History","Dashboards, full-text search, exports, and AI narratives",                      "5 capabilities"),
    ("🔀", "File Merge",         "Join separate Epic, Feature, and Story JIRA exports into one clean file",        "3 capabilities"),
    ("👥", "Teams",              "Team-level WAF analytics, cross-team epic views, and Team of Teams rollups",     "4 capabilities"),
    ("🏆", "Story Quality",      "AI-powered Definition of Ready scoring, chat refinement, and story rewrites",    "5 capabilities"),
    ("🏗",  "Epic Lineage",      "Epic health scoring, hierarchy trees, and cross-team epic assignment",           "3 capabilities"),
    ("⚡", "Disputes",           "Raise, queue, and resolve disagreements with AI classifications",                "3 capabilities"),
    ("⚙️", "Settings & Versions","WAF/GT file management, named versions, baselines, and system config",           "5 capabilities"),
]

rect(sl, Inches(0.55), Inches(1.86), Inches(12.2), Inches(0.36), fill=SURFACE, lc=BORDER, lw=Pt(0.5)).line.color.rgb = BORDER
txt(sl, "Module",        Inches(1.6),  Inches(1.92), Inches(2.5), Inches(0.24), size=9, bold=True, color=MUTED)
txt(sl, "Description",   Inches(4.1),  Inches(1.92), Inches(7.0), Inches(0.24), size=9, bold=True, color=MUTED)
txt(sl, "Scope",         Inches(11.4), Inches(1.92), Inches(1.2), Inches(0.24), size=9, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

for i, (ico, name, desc, count) in enumerate(modules_detail):
    ry = Inches(2.22) + i * Inches(0.6)
    bg_c = WHITE if i % 2 == 0 else SURFACE
    rect(sl, Inches(0.55), ry, Inches(12.2), Inches(0.56), fill=bg_c, lc=BORDER, lw=Pt(0.4))
    txt(sl, ico,   Inches(0.72), ry + Inches(0.1),  Inches(0.4),  Inches(0.4),  size=15)
    txt(sl, name,  Inches(1.22), ry + Inches(0.13), Inches(2.65), Inches(0.34), size=11.5, bold=True, color=HEAD)
    txt(sl, desc,  Inches(4.1),  ry + Inches(0.13), Inches(7.1),  Inches(0.34), size=10.5, color=TEXT)
    txt(sl, count, Inches(11.2), ry + Inches(0.13), Inches(1.35), Inches(0.34), size=11, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

footer(sl, 3)


# ════════════════════════════════════════════════════════════
#  SLIDE 4 — WHAT LEADERSHIP GETS  (MAIN)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Outcomes")
slide_title(sl, "What Leadership Gets")
slide_sub(sl, "The decisions and visibility this platform enables — without dictating the work that produces them.")
divider(sl, Inches(1.42))

# Top stat row — re-purposed to outcomes, not feature counts
stats = [
    ("Aligned",   "Work tied to priorities", "Every story mapped to WAF colour & category", ACCENT),
    ("Trusted",   "Auditable AI decisions",  "Reasoning, confidence, version logged",       LLM_FG),
    ("Consistent","One taxonomy, all teams", "Same rules applied at scale",                 ACCENT2),
    ("Actionable","Mismatches surfaced",     "Disagreements flagged and resolvable",        WARN),
]
sw = Inches(2.9)
for i, (val, lbl_t, sub, col) in enumerate(stats):
    sx = Inches(0.55) + i * (sw + Inches(0.18))
    rrect(sl, sx, Inches(1.6), sw, Inches(1.34), fill=WHITE, lc=BORDER)
    rect(sl, sx, Inches(1.6), sw, Pt(4), fill=col).line.fill.background()
    txt(sl, val,   sx, Inches(1.74), sw, Inches(0.5), size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, lbl_t, sx, Inches(2.26), sw, Inches(0.26), size=10.5, bold=True, color=HEAD, align=PP_ALIGN.CENTER)
    txt(sl, sub,   sx, Inches(2.54), sw, Inches(0.36), size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Decision-grade visibility list
txt(sl, "DECISION-GRADE VISIBILITY", Inches(0.55), Inches(3.18),
    Inches(6), Inches(0.26), size=9, bold=True, color=ACCENT)

visibility = [
    ("📊", "Spend by strategic priority",
     "Story count and proportion by WAF category and colour — per upload, per Team of Teams, per epic."),
    ("⚖️", "Run vs. Change split",
     "How much capacity is keep-the-lights-on vs. value-generating, sliced any way leadership asks."),
    ("🎯", "Mismatches flagged for review",
     "Where teams' tags disagree with AI — the conversations leadership needs to have, surfaced not buried."),
    ("🏆", "Story readiness signal",
     "What % of stories pass Definition of Ready before a PI starts — a quality leading indicator."),
    ("🔍", "Full audit trail",
     "Every classification linked to its reasoning, confidence, and the WAF version that produced it."),
    ("📤", "Stakeholder-ready exports",
     "One-click Excel and PDF for steering committees, audits, and quarterly reviews."),
]
vw, vh = Inches(3.9), Inches(1.5)
for i, (ico, ttl, dsc) in enumerate(visibility):
    vx = Inches(0.55) + (i % 3) * (vw + Inches(0.18))
    vy = Inches(3.5) + (i // 3) * (vh + Inches(0.14))
    rrect(sl, vx, vy, vw, vh, fill=WHITE, lc=BORDER)
    rect(sl, vx, vy, vw, Pt(3), fill=ACCENT).line.fill.background()
    txt(sl, ico, vx + Inches(0.14), vy + Inches(0.1), Inches(0.36), Inches(0.36), size=14)
    txt(sl, ttl, vx + Inches(0.55), vy + Inches(0.12), vw - Inches(0.65), Inches(0.3),
        size=11, bold=True, color=HEAD)
    txt(sl, dsc, vx + Inches(0.14), vy + Inches(0.5), vw - Inches(0.26), Inches(0.92),
        size=9.5, color=MUTED)

# Closing line
txt(sl, "WAF Classifier  ·  Built to bring consistency, auditability, and quality to every work item.",
    Inches(0.55), Inches(6.86), Inches(12.0), Inches(0.3),
    size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

footer(sl, 4)


# ════════════════════════════════════════════════════════════
#  SLIDE 5 — WHERE THE AI IS USED  (MAIN — governance / trust)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Governance")
slide_title(sl, "Where the AI is Used")
slide_sub(sl, "7 of 35 capabilities invoke the LLM. Everything else is deterministic — file parsing, database queries, rule-based logic.")
divider(sl, Inches(1.42))

txt(sl, "AI CALLS", Inches(0.55), Inches(1.5), Inches(4), Inches(0.24), size=8.5, bold=True, color=ACCENT)
llm_calls = [
    ("💬", "Single-story classification",  "WAF category, colour, run/change, confidence"),
    ("📦", "Bulk file verification",        "Same classification at scale with concurrent workers"),
    ("📝", "Executive narrative",           "AI-written paragraph summary of classification stats"),
    ("🎯", "DoR scoring",                   "9-criterion pass/fail scored per story"),
    ("💬", "Story improvement chat",        "Targeted fix suggestions per failing criterion"),
    ("✍️", "Story rewrite",                 "Full story rewrite addressing all failing criteria"),
]
for i, (ico, ttl, sub) in enumerate(llm_calls):
    ry = Inches(1.82) + i * Inches(0.82)
    rrect(sl, Inches(0.55), ry, Inches(5.9), Inches(0.72), fill=WHITE, lc=BORDER)
    txt(sl, ico, Inches(0.72), ry + Inches(0.1), Inches(0.38), Inches(0.5), size=14)
    txt(sl, ttl, Inches(1.22), ry + Inches(0.06), Inches(3.9), Inches(0.28), size=11, bold=True, color=HEAD)
    txt(sl, sub, Inches(1.22), ry + Inches(0.36), Inches(4.1), Inches(0.3),  size=9.5, color=MUTED)
    r = rrect(sl, Inches(5.72), ry + Inches(0.22), Inches(0.65), Pt(15), fill=LLM_BG, lc=LLM_BG)
    txt(sl, "🤖 LLM", Inches(5.7), ry + Inches(0.2), Inches(0.7), Pt(17), size=7.5, bold=True, color=LLM_FG, align=PP_ALIGN.CENTER)

txt(sl, "LLM USAGE BY MODULE", Inches(7.1), Inches(1.5), Inches(5.5), Inches(0.24), size=8.5, bold=True, color=ACCENT)
bars = [
    ("Classification",  6, 7),
    ("Analytics",       1, 5),
    ("File Merge",      0, 3),
    ("Teams",           0, 4),
    ("Story Quality",   3, 5),
    ("Epic Lineage",    0, 3),
    ("Disputes",        0, 3),
    ("Settings",        0, 5),
]
max_bar = Inches(3.8)
for i, (name, used, tot) in enumerate(bars):
    ry = Inches(1.82) + i * Inches(0.62)
    txt(sl, name, Inches(7.1), ry, Inches(1.6), Inches(0.36), size=10.5, bold=True, color=TEXT)
    track = rrect(sl, Inches(8.82), ry + Inches(0.08), max_bar, Inches(0.2), fill=SURFACE, lc=BORDER, lw=Pt(0.4))
    if used > 0:
        fw = int(max_bar * used / tot)
        r = rrect(sl, Inches(8.82), ry + Inches(0.08), fw, Inches(0.2), fill=ACCENT, lc=ACCENT, lw=Pt(0))
        r.line.fill.background()
    txt(sl, f"{used}/{tot}", Inches(12.72), ry, Inches(0.52), Inches(0.36),
        size=10, bold=(used > 0), color=ACCENT if used > 0 else MUTED, align=PP_ALIGN.RIGHT)

rrect(sl, Inches(7.1), Inches(6.72) - Inches(0.38), Inches(5.75), Inches(0.78), fill=SURFACE, lc=BORDER)
txt(sl, "GATEWAY SUPPORT", Inches(7.3), Inches(6.72) - Inches(0.34), Inches(3), Inches(0.22), size=8.5, bold=True, color=MUTED)
gws = ["Anthropic", "AWS Bedrock", "PortKey", "Apigee (OAuth2)"]
gx = Inches(7.3)
for gw in gws:
    gw_w = Inches(1.3)
    rrect(sl, gx, Inches(6.72) - Inches(0.06), gw_w, Inches(0.3), fill=LLM_BG, lc=LLM_BG)
    txt(sl, gw, gx, Inches(6.72) - Inches(0.07), gw_w, Inches(0.28), size=8.5, bold=True, color=LLM_FG, align=PP_ALIGN.CENTER)
    gx += gw_w + Inches(0.08)

footer(sl, 5)


# ════════════════════════════════════════════════════════════
#  SLIDE 6 — APPENDIX DIVIDER
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=MUTED).line.fill.background()

# Centred divider treatment
txt(sl, "APPENDIX",
    Inches(0.55), Inches(2.9), Inches(12.2), Inches(0.4),
    size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txt(sl, "Module Detail",
    Inches(0.55), Inches(3.32), Inches(12.2), Inches(0.9),
    size=44, bold=True, color=HEAD, align=PP_ALIGN.CENTER)

txt(sl, "Capability-by-capability deep dives across all 8 modules.",
    Inches(0.55), Inches(4.3), Inches(12.2), Inches(0.4),
    size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Mini index
index_items = [
    ("A1", "🤖 Classification"),
    ("A2", "📊 Analytics & History"),
    ("A3", "🔀 File Merge"),
    ("A4", "👥 Teams"),
    ("A5", "🏆 Story Quality"),
    ("A6", "🏗 Epic Lineage · ⚡ Disputes"),
    ("A7", "⚙️ Settings & Version Control"),
]
iw = Inches(11.0)
ix0 = (SLIDE_W - iw) / 2
iy0 = Inches(5.0)
rrect(sl, ix0, iy0, iw, Inches(1.7), fill=WHITE, lc=BORDER)
for i, (tag, label) in enumerate(index_items):
    col = i % 2
    row = i // 2
    cellx = ix0 + Inches(0.4) + col * Inches(5.2)
    celly = iy0 + Inches(0.2) + row * Inches(0.4)
    txt(sl, tag, cellx, celly, Inches(0.55), Inches(0.32),
        size=10, bold=True, color=ACCENT)
    txt(sl, label, cellx + Inches(0.7), celly, Inches(4.4), Inches(0.32),
        size=11, color=HEAD)

footer(sl, 6)


# ════════════════════════════════════════════════════════════
#  SLIDE 7 — A1 · CLASSIFICATION  (APPENDIX)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Appendix · A1 · Module 1")
slide_title(sl, "🤖 Classification")
slide_sub(sl, "Classify work items against the WAF — one at a time interactively, or thousands at once via async bulk processing.")
divider(sl, Inches(1.42))

cards = [
    ("💬", "Single-Story Classification",
     "Type a story title and description into a chat interface. The AI reads your WAF definitions and ground truth examples, then returns a category, colour, Run/Change flag, and confidence score — with reasoning you can challenge in follow-up messages.",
     True),
    ("📦", "Bulk File Verification",
     "Upload a CSV or Excel of stories. The tool auto-detects your columns, maps them, then classifies every row using concurrent AI workers. Results stream back as a progress bar — you review, approve, or dispute each row before saving.",
     True),
    ("🗂", "Column Mapping Preview",
     "Before bulk verify runs, the tool parses your file and suggests which column maps to which field (title, WAF category, Team of Teams, story points, etc.). You correct the mapping before AI kicks off — no tokens spent.",
     False),
    ("✅", "Approve to Ground Truth",
     "One-click promotion of any verified story into your ground truth library. Future AI calls use it as a calibration example, improving classification consistency over time as your dataset grows.",
     False),
]

cw, ch = Inches(5.95), Inches(2.38)
for i, (ico, ttl, dsc, llm) in enumerate(cards):
    cx = Inches(0.55) + (i % 2) * (cw + Inches(0.18))
    cy = Inches(1.55) + (i // 2) * (ch + Inches(0.14))
    cap_card(sl, cx, cy, cw, ch, ico, ttl, dsc, llm)

footer(sl, 7)


# ════════════════════════════════════════════════════════════
#  SLIDE 8 — A2 · ANALYTICS  (APPENDIX)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Appendix · A2 · Module 2")
slide_title(sl, "📊 Analytics & History")
slide_sub(sl, "Understand how your backlog is classified — across teams, time periods, and WAF categories — with flexible exports for stakeholders.")
divider(sl, Inches(1.42))

cards = [
    ("🏠", "Dashboard Summary",
     "At-a-glance totals: stories classified, approval rate, mismatch rate, top categories, colour distribution, and a daily trend line — filterable by team, upload, or date range.", False),
    ("🔍", "Story Drilldown",
     "Full-text searchable table of every classified story. Filter by team, category, colour, match status, and date. Expand any row to see AI reasoning, original tag, and full classification detail.", False),
    ("📝", "AI Executive Narrative",
     "Generates a concise, human-readable paragraph from your summary stats — describing WAF alignment, risk areas, and team performance. Ready to paste into a slide deck or status update.", True),
    ("📤", "CSV / Excel Export",
     "Download filtered history as a flat CSV or a formatted multi-sheet Excel workbook (Summary, Raw Data tabs) for offline analysis or stakeholder sharing.", False),
    ("⬆️", "Bulk Import",
     "Upload a pre-classified CSV/Excel to populate the history database without running AI — preserving timestamps and original tags from previous tools or manual tagging exercises.", False),
]

# 5 cards: top row 3, bottom row 2 (centred)
cw, ch = Inches(3.9), Inches(2.16)
for i in range(3):
    ico, ttl, dsc, llm = cards[i]
    cx = Inches(0.55) + i * (cw + Inches(0.18))
    cap_card(sl, cx, Inches(1.55), cw, ch, ico, ttl, dsc, llm)
# Bottom 2, centred under the 3-row
bot_total_w = 2 * cw + Inches(0.18)
bot_x0 = (SLIDE_W - bot_total_w) / 2
for i in range(2):
    ico, ttl, dsc, llm = cards[3 + i]
    cx = bot_x0 + i * (cw + Inches(0.18))
    cap_card(sl, cx, Inches(1.55) + ch + Inches(0.18), cw, ch, ico, ttl, dsc, llm)

footer(sl, 8)


# ════════════════════════════════════════════════════════════
#  SLIDE 9 — A3 · MERGE  (APPENDIX)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Appendix · A3 · Module 3")
slide_title(sl, "🔀 File Merge")
slide_sub(sl, "Join separate JIRA exports — Epics, Features, and Stories — into one clean, enriched file ready for classification or reporting.")
divider(sl, Inches(1.42))

cards = [
    ("🧩", "Epic → Feature → Story Join",
     "Takes up to three JIRA export files and joins them by name (case-insensitive). Extracts WAF colour and category, strips Run/Change suffixes from Epic Names, and propagates PI Name down to every story row.",
     False),
    ("⚠️", "Issue Detection",
     "Flags four problem types during the merge: orphaned stories (no matching feature), orphaned features (no matching epic), stories missing WAF data, and unrecognised WAF colours. You choose which rows to include before downloading.",
     False),
    ("🚀", "Send to Classifier",
     "One-click handoff from the merge result directly into the Bulk Verify flow — no intermediate download needed. The merged file arrives with column detection already run and all fields pre-mapped.",
     False),
]

cw, ch = Inches(3.9), Inches(2.06)
for i, (ico, ttl, dsc, llm) in enumerate(cards):
    cap_card(sl, Inches(0.55) + i * (cw + Inches(0.18)), Inches(1.55), cw, ch, ico, ttl, dsc, llm)

bx, by = Inches(0.55), Inches(3.82)
rrect(sl, bx, by, Inches(12.2), Inches(2.92), fill=WHITE, lc=BORDER)
txt(sl, "OUTPUT SCHEMA — 18 COLUMNS", bx + Inches(0.2), by + Inches(0.14),
    Inches(5), Inches(0.25), size=8.5, bold=True, color=MUTED)

cols_row1 = ["PI Name","Epic Id","Epic Name","Epic Desc","Block","WAF","WAF Color","WAF Category","Run/Change"]
cols_row2 = ["Feature Id","Feature Name","Feature Desc","Team of Teams","Story Id","Story Name","Story Desc","Story Points","Assigned Teams"]
waf_cols  = {"WAF","WAF Color","WAF Category","Run/Change"}
tot_cols  = {"Team of Teams"}

for row_i, cols in enumerate([cols_row1, cols_row2]):
    for j, col in enumerate(cols):
        cx = bx + Inches(0.18) + j * Inches(1.34)
        cy = by + Inches(0.46) + row_i * Inches(0.52)
        colf = ACCENT_PALE if col in waf_cols else TEAL_PALE if col in tot_cols else SURFACE
        colr = ACCENT if col in waf_cols else ACCENT2 if col in tot_cols else HEAD
        rrect(sl, cx, cy, Inches(1.25), Inches(0.38), fill=colf, lc=BORDER, lw=Pt(0.4))
        txt(sl, col, cx, cy + Inches(0.03), Inches(1.25), Inches(0.32),
            size=8.5, bold=(col in waf_cols or col in tot_cols), color=colr, align=PP_ALIGN.CENTER)

lx = bx + Inches(0.2)
ly = by + Inches(1.62)
for lbl, bg_c, fg_c in [("WAF fields", ACCENT_PALE, ACCENT), ("Team of Teams", TEAL_PALE, ACCENT2)]:
    rrect(sl, lx, ly, Inches(0.14), Inches(0.14), fill=bg_c, lc=bg_c)
    txt(sl, lbl, lx + Inches(0.2), ly - Pt(1), Inches(1.5), Inches(0.18), size=8, color=MUTED)
    lx += Inches(1.9)

footer(sl, 9)


# ════════════════════════════════════════════════════════════
#  SLIDE 10 — A4 · TEAMS  (APPENDIX)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Appendix · A4 · Module 4")
slide_title(sl, "👥 Teams")
slide_sub(sl, "Understand how every Scrum Team and Team of Teams (ART) is aligned to the WAF — and how their work distributes across epics.")
divider(sl, Inches(1.42))

cards = [
    ("📋", "Team Summary",
     "For every Scrum Team: total stories, unique epics touched, WAF category breakdown, dominant colour, mismatch rate, and their Team of Teams (ART). Supports ART-level rollup filtering.", False),
    ("🔎", "Team Drilldown",
     "Drill into a single team to see their work organised as Epic → Feature → Story trees, with per-epic category consistency, dominant classifications, and health indicators.", False),
    ("🗺", "Teams by Epic",
     "Pick an epic and see every team contributing to it, their story counts, and category distribution. Immediately surfaces cross-team alignment (or divergence) on a single initiative.", False),
    ("📜", "Epics List",
     "A master list of all epics with story counts, team assignments, and dominant WAF category — the navigation starting point for the By Epic tab and cross-team analysis.", False),
]

cw, ch = Inches(5.95), Inches(2.16)
for i, (ico, ttl, dsc, llm) in enumerate(cards):
    cx = Inches(0.55) + (i % 2) * (cw + Inches(0.18))
    cy = Inches(1.55) + (i // 2) * (ch + Inches(0.14))
    cap_card(sl, cx, cy, cw, ch, ico, ttl, dsc, llm)

rrect(sl, Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.82), fill=SURFACE, lc=BORDER)
txt(sl, "HIERARCHY", Inches(0.75), Inches(6.3), Inches(1.2), Inches(0.24), size=8.5, bold=True, color=MUTED)
nodes = [
    ("Team of Teams (ART)", ACCENT_PALE, ACCENT),
    ("→", None, MUTED),
    ("Scrum Teams", TEAL_PALE, ACCENT2),
    ("→", None, MUTED),
    ("Stories", SURFACE, HEAD),
]
nx = Inches(2.2)
for (node_t, bg_c, fg_c) in nodes:
    if node_t == "→":
        txt(sl, "→", nx, Inches(6.36), Inches(0.34), Inches(0.36), size=14, color=MUTED, align=PP_ALIGN.CENTER)
        nx += Inches(0.36)
    else:
        nw = Inches(2.3)
        rrect(sl, nx, Inches(6.3), nw, Inches(0.44), fill=bg_c, lc=BORDER, lw=Pt(0.5))
        txt(sl, node_t, nx, Inches(6.34), nw, Inches(0.36), size=11, bold=True, color=fg_c, align=PP_ALIGN.CENTER)
        nx += nw + Inches(0.14)

footer(sl, 10)


# ════════════════════════════════════════════════════════════
#  SLIDE 11 — A5 · STORY QUALITY  (APPENDIX)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Appendix · A5 · Module 5")
slide_title(sl, "🏆 Story Quality — Definition of Ready")
slide_sub(sl, "Score every story against a 9-criterion DoR rubric — then use AI to refine or fully rewrite stories that don't pass.")
divider(sl, Inches(1.42))

cards_top = [
    ("🎯", "Batch DoR Scoring",
     "Scores every story against 9 criteria: Narrative Format, Source Data, Business Rules, Output Artifact, Acceptance Criteria, Data Quality Checks, Traceability, Story Points, Dependencies. Each criterion gets a pass/fail with a short reason.",
     True),
    ("💬", "Story Improvement Chat",
     "Open a chat against a specific story. The AI suggests targeted improvements for whichever criteria are failing — collaborative, iterative refinement without rewriting the whole story.",
     True),
    ("✍️", "AI Story Rewrite",
     "One click for the AI to produce a fully rewritten version of the story addressing all failing criteria, in the correct narrative format — ready to copy back into JIRA.",
     True),
]
cards_bot = [
    ("📊", "Quality Results & History",
     "Browse scoring results by upload, team, or scoring run. See aggregate Ready / Needs Work / Not Ready counts, per-criterion pass rates, and score trends across multiple runs.",
     False),
    ("📤", "Quality Export",
     "Download results as CSV — one row per story with pass/fail for all 9 criteria, overall score percentage, and team/epic metadata for downstream reporting.",
     False),
]

cw_t, ch_t = Inches(3.9), Inches(1.9)
for i, (ico, ttl, dsc, llm) in enumerate(cards_top):
    cap_card(sl, Inches(0.55) + i * (cw_t + Inches(0.18)), Inches(1.55), cw_t, ch_t, ico, ttl, dsc, llm)

cw_b, ch_b = Inches(5.95), Inches(1.76)
for i, (ico, ttl, dsc, llm) in enumerate(cards_bot):
    cap_card(sl, Inches(0.55) + i * (cw_b + Inches(0.18)), Inches(3.62), cw_b, ch_b, ico, ttl, dsc, llm)

bx, by = Inches(0.55), Inches(5.52)
rrect(sl, bx, by, Inches(12.2), Inches(1.16), fill=WHITE, lc=BORDER)
txt(sl, "SCORE BANDS", bx + Inches(0.2), by + Inches(0.1), Inches(2), Inches(0.24), size=8.5, bold=True, color=MUTED)
bands = [("✅  Ready", "≥ 89%", GREEN), ("⚠️  Needs Work", "56 – 88%", WARN), ("❌  Not Ready", "< 56%", RED)]
for j, (name, pct, col) in enumerate(bands):
    bsx = bx + Inches(0.5) + j * Inches(3.8)
    rrect(sl, bsx, by + Inches(0.44), Inches(3.4), Inches(0.56), fill=SURFACE, lc=BORDER, lw=Pt(0.4))
    txt(sl, name, bsx + Inches(0.14), by + Inches(0.5),  Inches(2.2), Inches(0.3), size=11, bold=True, color=col)
    txt(sl, pct,  bsx + Inches(2.5),  by + Inches(0.5),  Inches(0.8), Inches(0.3), size=11, color=MUTED, align=PP_ALIGN.RIGHT)

footer(sl, 11)


# ════════════════════════════════════════════════════════════
#  SLIDE 12 — A6 · EPIC LINEAGE + DISPUTES  (APPENDIX)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Appendix · A6 · Modules 6 & 7")
slide_title(sl, "🏗  Epic Lineage   ·   ⚡ Disputes")
divider(sl, Inches(1.0))

txt(sl, "EPIC LINEAGE", Inches(0.55), Inches(1.08), Inches(5), Inches(0.24), size=8.5, bold=True, color=ACCENT)
txt(sl, "DISPUTES",     Inches(6.88), Inches(1.08), Inches(5), Inches(0.24), size=8.5, bold=True, color=ACCENT)
rect(sl, Inches(6.65), Inches(1.06), Pt(1), Inches(5.8), fill=BORDER).line.fill.background()

lineage = [
    ("❤️", "Epic Health Scoring",
     "Each epic gets a 0–100 health score based on: colour consistency (40%), WAF category focus (30%), distinct colour count (20%), and mismatch rate (10%). Epics spanning 3+ colours are flagged as 'mixed'.", False),
    ("🌳", "Hierarchy Tree",
     "Full Epic → Feature → Story tree showing which teams deliver each part, story counts per node, and dominant classifications at every level of the hierarchy.", False),
    ("🔗", "Bulk Epic Assignment",
     "Assign or re-assign Epic and Feature fields across multiple stories at once — essential for cleaning up imported data where those fields were blank or inconsistent.", False),
]
disputes_data = [
    ("🚩", "Raise a Dispute",
     "Flag any AI classification you disagree with. Submit your suggested category and a comment — the original AI reasoning is preserved alongside your correction for reviewer context.", False),
    ("📋", "Dispute Queue",
     "Filtered list of all disputes by status: Pending, Accepted, Dismissed, or Flagged for WAF Review. Shows AI answer vs. suggested answer with reviewer notes.", False),
    ("⚖️", "Resolve a Dispute",
     "Three paths: Dismiss (AI was right), Accept (save correction + optionally promote to ground truth), or Flag for WAF Review (the WAF definition itself may need updating).", False),
]

cw, ch = Inches(5.75), Inches(1.68)
for i, (ico, ttl, dsc, llm) in enumerate(lineage):
    cap_card(sl, Inches(0.55), Inches(1.38) + i * (ch + Inches(0.12)), cw, ch, ico, ttl, dsc, llm)
for i, (ico, ttl, dsc, llm) in enumerate(disputes_data):
    cap_card(sl, Inches(6.88), Inches(1.38) + i * (ch + Inches(0.12)), cw, ch, ico, ttl, dsc, llm)

footer(sl, 12)


# ════════════════════════════════════════════════════════════
#  SLIDE 13 — A7 · SETTINGS  (APPENDIX)
# ════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, Inches(0.06), SLIDE_H, fill=ACCENT).line.fill.background()
slide_eyebrow(sl, "Appendix · A7 · Module 8")
slide_title(sl, "⚙️  Settings & Version Control")
slide_sub(sl, "Manage the AI's knowledge base — WAF definitions and ground truth — with full version history, baselines, and system configuration.")
divider(sl, Inches(1.42))

cards = [
    ("📖", "WAF Definitions Management",
     "Upload a new WAF definitions file (CSV, Excel, JSON, or plain text) at any time. The tool auto-parses categories, colours, and descriptions. Edit definitions inline from the Settings page without re-uploading.", False),
    ("🎓", "Ground Truth Management",
     "Upload, add, edit, or delete ground truth examples — the calibration examples sent to the AI in every prompt. Quality of ground truth directly drives classification consistency across all teams and uploads.", False),
    ("🏷", "Named Versions",
     "Save the current WAF definitions or ground truth as a named, timestamped version (e.g. 'Q3 Baseline'). Switch freely — bulk verify jobs record which version was active, ensuring full reproducibility.", False),
    ("📸", "Baseline Snapshots",
     "Save a paired snapshot of WAF + ground truth as a named baseline. Restore in one click — ideal for audit trails, rolling back after an experiment, or matching a prior quarter's configuration.", False),
]

cw, ch = Inches(5.95), Inches(2.2)
for i, (ico, ttl, dsc, llm) in enumerate(cards):
    cx = Inches(0.55) + (i % 2) * (cw + Inches(0.18))
    cy = Inches(1.55) + (i // 2) * (ch + Inches(0.14))
    cap_card(sl, cx, cy, cw, ch, ico, ttl, dsc, llm)

rrect(sl, Inches(0.55), Inches(6.22), Inches(12.2), Inches(0.84), fill=SURFACE, lc=BORDER)
txt(sl, "SYSTEM CONFIG — NO CODE NEEDED", Inches(0.75), Inches(6.3), Inches(4), Inches(0.24), size=8.5, bold=True, color=MUTED)
configs = [("Batch size", "stories per AI call"), ("Workers", "concurrent AI threads"),
           ("Rate limit", "calls per minute"), ("Gateway", "Anthropic · Bedrock · PortKey · Apigee")]
for j, (k, v) in enumerate(configs):
    cx = Inches(0.75) + j * Inches(3.0)
    txt(sl, k, cx, Inches(6.3),  Inches(2.8), Inches(0.26), size=10.5, bold=True, color=ACCENT)
    txt(sl, v, cx, Inches(6.56), Inches(2.8), Inches(0.26), size=9.5, color=MUTED)

footer(sl, 13)


# ── Save ─────────────────────────────────────────────────────
out = "static/WAF_Classifier_Capabilities.pptx"
prs.save(out)
print(f"✓  Saved → {out}  (13 slides: 5 main · 1 divider · 7 appendix)")
